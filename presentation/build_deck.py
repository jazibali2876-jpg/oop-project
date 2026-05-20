"""Generate RestaurantPOS_Demo.pptx for the Smart Restaurant POS final-year demo.

Grounded in the planning docs and built code. No invented features.
Run from the project root:  python presentation/build_deck.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# -- Palette (matches the in-app theme tokens) ------------------------------
BRAND   = RGBColor(0xFF, 0x6B, 0x35)
INK     = RGBColor(0x0F, 0x11, 0x15)
SOFT    = RGBColor(0x17, 0x1A, 0x21)
LIGHT   = RGBColor(0xF3, 0xF4, 0xF6)
MUTED   = RGBColor(0x9C, 0xA3, 0xAF)
LINE    = RGBColor(0x26, 0x2B, 0x33)
SUCCESS = RGBColor(0x22, 0xC5, 0x5E)
INFO    = RGBColor(0x3B, 0x82, 0xF6)
WARN    = RGBColor(0xF5, 0x9E, 0x0B)

# -- Slide geometry: 16:9 -----------------------------------------------------
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

ROOT = Path(__file__).resolve().parent.parent
SHOTS_DIR = ROOT / "presentation" / "screenshots"


# -- Helpers ------------------------------------------------------------------
def add_blank():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, rgb):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill(s, rgb)
    return s


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=LIGHT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_paragraphs(slide, x, y, w, h, lines, *, size=14, color=LIGHT, bullet=True, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ("•  " + ln) if bullet else ln
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def slide_chrome(slide, title, subtitle=None, dark=True):
    bg = SOFT if dark else LIGHT
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, bg)
    # Brand accent strip on the left
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, BRAND)
    title_color = LIGHT if dark else INK
    add_text(slide, Inches(0.55), Inches(0.35), Inches(12.3), Inches(0.7),
             title, size=30, bold=True, color=title_color)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(1.05), Inches(12.3), Inches(0.5),
                 subtitle, size=14, color=MUTED if dark else MUTED)
    # Brand underline
    add_rect(slide, Inches(0.55), Inches(1.55), Inches(0.9), Inches(0.04), BRAND)


def add_screenshot_or_placeholder(slide, name, x, y, w, h):
    """Insert <name>.png from presentation/screenshots/, else a labeled placeholder."""
    p = SHOTS_DIR / f"{name}.png"
    if p.exists():
        slide.shapes.add_picture(str(p), x, y, w, h)
        return True
    box = add_rect(slide, x, y, w, h, LINE)
    box.line.color.rgb = BRAND
    box.line.width = Pt(1)
    add_text(slide, x, y + Inches(0.1), w, Inches(0.4),
             f"[screenshot: {name}]", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + h - Inches(0.4), w, Inches(0.3),
             "(capture and drop into presentation/screenshots/)",
             size=10, color=MUTED, align=PP_ALIGN.CENTER)
    return False


# ============================================================================
# Slides
# ============================================================================

# 1. Cover
s = add_blank()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, INK)
add_rect(s, 0, Inches(3.45), SLIDE_W, Inches(0.08), BRAND)
add_text(s, Inches(0.7), Inches(1.6), Inches(12), Inches(1.5),
         "Smart Restaurant POS", size=58, bold=True, color=LIGHT)
add_text(s, Inches(0.7), Inches(2.5), Inches(12), Inches(0.5),
         "A modern point-of-sale desktop application", size=22, color=BRAND)
add_text(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.5),
         "C++17  ·  SFML 3.1  ·  Dear ImGui  ·  CMake  ·  MinGW-w64",
         size=18, color=MUTED)
add_text(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.4),
         "Final-year Object-Oriented Programming Project", size=16, color=LIGHT)
add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.4),
         "Jazib Ali", size=16, color=LIGHT)
add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
         "2026", size=14, color=MUTED)


# 2. Problem & objective
s = add_blank(); slide_chrome(s, "Problem & Objective")
add_text(s, Inches(0.55), Inches(1.9), Inches(12), Inches(0.5),
         "The problem", size=20, bold=True, color=BRAND)
add_paragraphs(s, Inches(0.55), Inches(2.4), Inches(12), Inches(2.0), [
    "Independent restaurants juggle paper order pads, separate kitchen tickets, and standalone receipt printers.",
    "Inventory drifts out of sync. Reporting is manual. Customer loyalty data is lost.",
    "Commercial POS suites (Toast, Square) solve this, but are expensive, cloud-bound, and overkill for one location.",
], size=15)
add_text(s, Inches(0.55), Inches(4.6), Inches(12), Inches(0.5),
         "Our objective", size=20, bold=True, color=BRAND)
add_paragraphs(s, Inches(0.55), Inches(5.1), Inches(12), Inches(2.0), [
    "Build an offline desktop POS that looks and feels like a real product — covering orders, billing, inventory, tables, customers, kitchen, and analytics.",
    "Use the app to demonstrate strong object-oriented programming and complete file-handling, end-to-end.",
    "Stay realistically buildable inside one semester — no databases, no cloud, no enterprise frameworks.",
], size=15)


# 3. Tech stack
s = add_blank(); slide_chrome(s, "Tech Stack", "Everything is vendored — the app runs offline after a single clone.")
stack_rows = [
    ("Language",       "C++17"),
    ("UI framework",   "Dear ImGui (immediate-mode, header-light)"),
    ("Windowing",      "SFML 3.1 (window, events, audio, texture loading)"),
    ("Rendering",      "ImGui_ImplOpenGL3 backend (OpenGL 3.3)"),
    ("Toolchain",      "MinGW-w64 GCC 14.2 (UCRT / POSIX / SEH)"),
    ("Build system",   "CMake 3.30 (vendored portable)"),
    ("Storage",        "Plain C++ file handling — text + binary; NO database"),
    ("Charts",         "Hand-drawn with ImDrawList primitives — no ImPlot"),
]
y = Inches(2.0)
for label, value in stack_rows:
    add_text(s, Inches(0.7),  y, Inches(2.6), Inches(0.4), label,
             size=14, color=MUTED)
    add_text(s, Inches(3.4),  y, Inches(9.0), Inches(0.4), value,
             size=15, color=LIGHT, bold=True)
    y += Inches(0.46)


# 4. System architecture
s = add_blank(); slide_chrome(s, "System Architecture",
                              "Two layers + a manual integration glue layer.")
# UI layer block
add_rect(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(1.55), LINE)
add_text(s, Inches(0.8), Inches(2.05), Inches(11), Inches(0.4),
         "UI layer  ·  src/ui/",  size=16, bold=True, color=BRAND)
add_paragraphs(s, Inches(0.8), Inches(2.45), Inches(11.7), Inches(1.0), [
    "Immediate-mode screens (Splash, Login, Dashboard, Menu, Order, Billing, Inventory, Tables, Kitchen, Analytics).",
    "Holds only transient state (selection, scroll, animation timers). Never opens a file directly.",
], size=12)

# Arrow down
add_text(s, Inches(6.3), Inches(3.55), Inches(1), Inches(0.4), "↓ calls into",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Domain layer block
add_rect(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(1.55), LINE)
add_text(s, Inches(0.8), Inches(4.1), Inches(11), Inches(0.4),
         "Domain layer  ·  src/domain/  ·  no ImGui, no SFML",
         size=16, bold=True, color=SUCCESS)
add_paragraphs(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.0), [
    "Entities + Services + Repositories. Holds ALL business state; this is the layer graded for OOP.",
    "Services orchestrate use cases; repositories are the only code that opens fstreams.",
], size=12)

add_text(s, Inches(6.3), Inches(5.6), Inches(1), Inches(0.4), "↓ read / write",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Files block
add_rect(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(1.0), LINE)
add_text(s, Inches(0.8), Inches(6.15), Inches(11), Inches(0.4),
         "Data files  ·  data/  ·  text (.txt) + binary (.dat)",
         size=15, bold=True, color=INFO)
add_text(s, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.4),
         "users.dat   menu.txt   orders.dat   customers.dat   inventory.txt   tables.txt   sales_history.dat   receipts/*",
         size=11, color=MUTED)

# Side note: integration
add_text(s, Inches(0.55), Inches(7.2), Inches(13), Inches(0.3),
         "Glue:  src/platform/  —  hand-written ImGui ↔ SFML 3 integration over the bundled imgui_impl_opengl3 backend.",
         size=11, color=BRAND)


# 5. OOP demonstrated
s = add_blank(); slide_chrome(s, "OOP Demonstrated",
                              "Each concept maps to a concrete class in the built code.")
oop = [
    ("Classes & objects",       "MenuItem, Order, Customer, Ingredient, Table"),
    ("Encapsulation",           "Order — private items vector, mutation only via add / removeAt / setStatus"),
    ("Inheritance",             "User  →  Admin / Cashier / Kitchen   in  src/domain/auth/"),
    ("Polymorphism",            "Session::user().can(Capability) — runtime type decides"),
    ("Abstraction",             "User is abstract; pure-virtual can() and roleName()"),
    ("Virtual functions",       "User::can, ::roleName, ::role — overridden by every concrete role"),
    ("Templates",               "BinaryRepository<T> instantiated for User/Customer/Loyalty/Sales records;  Id<Tag>"),
    ("Operator overloading",    "Money +/-/*/<<;  Order::operator+=;  Receipt::operator<< (thermal layout)"),
    ("Exception handling",      "DomainException base + 8 subclasses; services throw, screens catch → toasts"),
    ("STL containers",          "vector, map, unordered_map, array<int,24>, optional, deque, unique_ptr"),
    ("Modular programming",     "domain/ split into 10 sub-packages; separate ui/ and util/ layers"),
    ("Header / source split",   "Every non-template class has a matching .h and .cpp"),
]
y = Inches(1.95)
for label, value in oop:
    add_text(s, Inches(0.7),  y, Inches(3.0), Inches(0.32), label,
             size=12, color=MUTED)
    add_text(s, Inches(3.7),  y, Inches(9.4), Inches(0.32), value,
             size=12, color=LIGHT, bold=True)
    y += Inches(0.42)


# 6. File-handling design
s = add_blank(); slide_chrome(s, "File-Handling Design",
                              "C++ fstreams only — text files for editable data, binary records for ledgers.")
# Two columns
add_text(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(0.4),
         "Text files (.txt)", size=18, bold=True, color=BRAND)
add_paragraphs(s, Inches(0.6), Inches(2.5), Inches(6.0), Inches(3.5), [
    "menu.txt — id | name | category | priceCents | imagePath | available",
    "inventory.txt — id | name | unit | stock | threshold | costCents",
    "recipes.txt — menuItemId | ingId:qty,ingId:qty,...",
    "tables.txt — id | seats | status | reservedFor | reservedAt",
    "settings.txt — k=v (theme, rememberedUser, ...)",
    "receipts/RCPT-NNNNNN.txt — thermal-style printable receipts",
], size=11)

add_text(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(0.4),
         "Binary records (.dat)", size=18, bold=True, color=INFO)
add_paragraphs(s, Inches(6.9), Inches(2.5), Inches(6.0), Inches(3.5), [
    "users.dat — UserRecord (128 B, static_assert-locked)",
    "orders.dat — variable-length: OrderHeader (128 B) + N OrderItems (56 B)",
    "customers.dat — CustomerRecord (96 B)",
    "loyalty.dat — LoyaltyRecord (32 B)  append-only ledger",
    "sales_history.dat — SalesRecord (64 B)  flat scan for analytics",
], size=11)

add_text(s, Inches(0.6), Inches(5.9), Inches(12.5), Inches(0.4),
         "Crash safety", size=16, bold=True, color=SUCCESS)
add_paragraphs(s, Inches(0.6), Inches(6.3), Inches(12.5), Inches(1.1), [
    "Whole-file rewrites use write-temp-then-rename (util::writeAtomic) — power loss can't corrupt the live file.",
    "Binary appends rely on the OS append guarantee. Every POD record struct is static_assert-locked at its on-disk byte count.",
    "Missing files at boot = empty data. Default seed data is created the first time the app starts.",
], size=12)


# 7. Features (1/3) — Auth + Dashboard + Menu
s = add_blank(); slide_chrome(s, "Feature Walkthrough  (1 / 3)",
                              "Auth, dashboard, menu management.")
add_screenshot_or_placeholder(s, "login",     Inches(0.6),  Inches(2.0), Inches(4.0), Inches(2.5))
add_screenshot_or_placeholder(s, "dashboard", Inches(4.75), Inches(2.0), Inches(4.0), Inches(2.5))
add_screenshot_or_placeholder(s, "menu",      Inches(8.9),  Inches(2.0), Inches(4.0), Inches(2.5))
add_text(s, Inches(0.6),  Inches(4.55), Inches(4.0), Inches(0.4),
         "Login — 3 roles", size=13, bold=True, color=LIGHT)
add_text(s, Inches(0.6),  Inches(4.92), Inches(4.0), Inches(2.4),
         "Admin / Cashier / Kitchen — capability\nchecks gate every action. Salted hash;\n'remember me' persists username in settings.",
         size=11, color=MUTED)
add_text(s, Inches(4.75), Inches(4.55), Inches(4.0), Inches(0.4),
         "Dashboard", size=13, bold=True, color=LIGHT)
add_text(s, Inches(4.75), Inches(4.92), Inches(4.0), Inches(2.4),
         "Today's revenue, active orders, low-stock\ncount, top item. 7-day revenue line chart\ndrawn manually with ImDrawList.",
         size=11, color=MUTED)
add_text(s, Inches(8.9),  Inches(4.55), Inches(4.0), Inches(0.4),
         "Menu Management", size=13, bold=True, color=LIGHT)
add_text(s, Inches(8.9),  Inches(4.92), Inches(4.0), Inches(2.4),
         "Admin-only CRUD with category filter and\nsearch. Cashiers see the same screen as\nread-only by capability check.",
         size=11, color=MUTED)


# 8. Features (2/3) — Order + Billing + Inventory
s = add_blank(); slide_chrome(s, "Feature Walkthrough  (2 / 3)",
                              "Orders, billing, inventory.")
add_screenshot_or_placeholder(s, "order",     Inches(0.6),  Inches(2.0), Inches(4.0), Inches(2.5))
add_screenshot_or_placeholder(s, "billing",   Inches(4.75), Inches(2.0), Inches(4.0), Inches(2.5))
add_screenshot_or_placeholder(s, "inventory", Inches(8.9),  Inches(2.0), Inches(4.0), Inches(2.5))
add_text(s, Inches(0.6),  Inches(4.55), Inches(4.0), Inches(0.4),
         "New Order", size=13, bold=True, color=LIGHT)
add_text(s, Inches(0.6),  Inches(4.92), Inches(4.0), Inches(2.4),
         "Menu grid + cart with auto totals + tax +\nprep estimate. Customer lookup by phone\n(register inline). Discount + notes.",
         size=11, color=MUTED)
add_text(s, Inches(4.75), Inches(4.55), Inches(4.0), Inches(0.4),
         "Billing & Receipts", size=13, bold=True, color=LIGHT)
add_text(s, Inches(4.75), Inches(4.92), Inches(4.0), Inches(2.4),
         "Cash / card / online-sim. Itemized receipt\nin a thermal-style 44-column layout via\nReceipt::operator<<. Save + PDF-style export.",
         size=11, color=MUTED)
add_text(s, Inches(8.9),  Inches(4.55), Inches(4.0), Inches(0.4),
         "Inventory", size=13, bold=True, color=LIGHT)
add_text(s, Inches(8.9),  Inches(4.92), Inches(4.0), Inches(2.4),
         "Stock table with low-stock row highlight.\nAuto-deduction when an order is placed,\nvia recipes mapping items to ingredients.",
         size=11, color=MUTED)


# 9. Features (3/3) — Tables + Kitchen + Analytics
s = add_blank(); slide_chrome(s, "Feature Walkthrough  (3 / 3)",
                              "Tables, kitchen display, analytics with hand-drawn charts.")
add_screenshot_or_placeholder(s, "tables",    Inches(0.6),  Inches(2.0), Inches(4.0), Inches(2.5))
add_screenshot_or_placeholder(s, "kitchen",   Inches(4.75), Inches(2.0), Inches(4.0), Inches(2.5))
add_screenshot_or_placeholder(s, "analytics", Inches(8.9),  Inches(2.0), Inches(4.0), Inches(2.5))
add_text(s, Inches(0.6),  Inches(4.55), Inches(4.0), Inches(0.4),
         "Tables", size=13, bold=True, color=LIGHT)
add_text(s, Inches(0.6),  Inches(4.92), Inches(4.0), Inches(2.4),
         "Color-coded grid: Free / Occupied /\nReserved. Tap to reserve, occupy, or free.\nReservations attach a Customer record.",
         size=11, color=MUTED)
add_text(s, Inches(4.75), Inches(4.55), Inches(4.0), Inches(0.4),
         "Kitchen Display", size=13, bold=True, color=LIGHT)
add_text(s, Inches(4.75), Inches(4.92), Inches(4.0), Inches(2.4),
         "Three lanes (Pending / Preparing /\nReady). Live elapsed-time per ticket;\nkitchen role advances status with one tap.",
         size=11, color=MUTED)
add_text(s, Inches(8.9),  Inches(4.55), Inches(4.0), Inches(0.4),
         "Analytics", size=13, bold=True, color=LIGHT)
add_text(s, Inches(8.9),  Inches(4.92), Inches(4.0), Inches(2.4),
         "Today / week / month ranges. Bar, pie, and\nline charts — all drawn manually with\nImDrawList (no ImPlot dependency).",
         size=11, color=MUTED)


# 10. Smart / optional features
s = add_blank(); slide_chrome(s, "Smart & Optional Features",
                              "Built vs deferred — honest scope.")
add_text(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(0.4),
         "Implemented", size=18, bold=True, color=SUCCESS)
add_paragraphs(s, Inches(0.6), Inches(2.5), Inches(6.0), Inches(4.0), [
    "Popular-item detection — Report::topItems via SalesHistoryRepository.",
    "Low-stock alerts on dashboard, with row-tinted highlights on Inventory.",
    "Loyalty points + tier (New / Bronze / Silver / Gold) via LoyaltyRule.",
    "Dark / Light theme switch persisted in settings.txt.",
    "Splash screen with progress bar; per-frame easing animations.",
    "Auto-seed data on first run so a clean clone boots straight into a usable demo.",
    "Atomic file writes (write-temp-then-rename) for crash-safe persistence.",
], size=12)

add_text(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(0.4),
         "Deferred (stretch goals)", size=18, bold=True, color=WARN)
add_paragraphs(s, Inches(6.9), Inches(2.5), Inches(6.0), Inches(4.0), [
    "Predictive low-stock (days-until-out estimate) — heuristic only.",
    "Sound effects via SFML audio.",
    "Keyboard shortcuts (Ctrl+N, Ctrl+P, ...).",
    "QR-code table menus and face-recognition login (out of scope).",
    "Real PDF export — current export is text-based.",
    "Backup / restore commands from the UI.",
], size=12)


# 11. Challenges faced
s = add_blank(); slide_chrome(s, "Challenges Faced", "Real engineering, not invented drama.")
challenges = [
    ("SFML MinGW linkage",
     "The vendored SFML-3.1.0 is an MSVC build (.lib + .pdb). MinGW could not link it due to MSVC ↔ Itanium C++ ABI incompatibility. "
     "Fix: rebuilt SFML 3.1.0 from source with the bundled GCC into third_party/SFML-mingw/ (proper lib*.a archives + DLLs)."),
    ("Immediate-mode UI vs strong OOP",
     "Dear ImGui has no persistent widget objects — easy to bleed business logic into draw functions. "
     "Fix: hard split into a UI layer (transient) and a Domain layer (no ImGui include). Domain owns all state; UI calls into it per frame."),
    ("Manual ImGui ↔ SFML 3 integration",
     "There is no ImGui-SFML binding for SFML 3. Built our own InputBridge + ImGuiBackend "
     "translating sf::Event variants (event.getIf<...>()) into ImGuiIO events, driving imgui_impl_opengl3 against SFML's GL context."),
    ("Binary record layout",
     "POD records crossing on-disk and in-memory must keep stable sizes. "
     "Fix: #pragma pack(1) + static_assert(sizeof(R) == N) for every record (UserRecord = 128, OrderHeader = 128, OrderItem = 56, ...)."),
    ("First-run UX",
     "A blank install must not crash if data files don't exist. "
     "Repositories tolerate missing files; services seed sensible defaults (admin user, 5 menu items, 10 tables, ingredients, recipes)."),
]
y = Inches(1.95)
for title, body in challenges:
    add_text(s, Inches(0.6), y, Inches(12.2), Inches(0.34), title,
             size=14, bold=True, color=BRAND)
    add_text(s, Inches(0.6), y + Inches(0.36), Inches(12.2), Inches(0.7), body,
             size=11, color=LIGHT)
    y += Inches(1.05)


# 12. Testing approach
s = add_blank(); slide_chrome(s, "Testing Approach")
add_paragraphs(s, Inches(0.55), Inches(2.0), Inches(12), Inches(5), [
    "Compile-time checks: every binary record has static_assert(sizeof(R) == N) and static_assert(std::is_trivially_copyable_v<R>).",
    "Build verification: smart_pos.exe linked at every package boundary during Phase 2 (3 rounds — common+auth, business core, services).",
    "Preflight script (tools/preflight_check.ps1) validates compiler, SFML link, every runtime DLL beside the .exe, fonts on disk, and asset/data paths.",
    "Fresh-machine simulation: dist/ copied to %TEMP%, PATH stripped to only C:\\Windows\\System32, then smart_pos.exe launched — confirms zero hidden dependencies.",
    "Smoke test (Phase 1 Stage 1): minimal ImGui window inside SFML window before any feature work — guaranteed the manual integration is sound.",
    "Manual exercise of every screen with first-run + populated states: login, place order, deduct stock, advance kitchen ticket, save receipt, view analytics.",
    "Logging: every service-level error → data/log.txt with ISO8601 timestamps, available for post-run review.",
], size=12)


# 13. Future improvements
s = add_blank(); slide_chrome(s, "Future Improvements")
add_paragraphs(s, Inches(0.55), Inches(2.0), Inches(12), Inches(5), [
    "Real PDF export via a small library (e.g. libharu), instead of text-formatted .pdf.txt files.",
    "True thermal-printer driver hookup (ESC/POS over USB) for printing receipts.",
    "Migration / format-versioning header on binary files so on-disk schema can evolve safely.",
    "Backup-and-restore actions in the Settings screen (zip the data/ folder).",
    "Predictive low-stock alerts using a simple rolling average over sales_history.dat.",
    "Multi-station mode: file-watcher syncing orders between front-of-house and kitchen displays on a LAN share.",
    "Internationalization — externalize all UI strings to data/i18n/<lang>.txt and add a language switch.",
    "Accessibility pass: keyboard-only navigation, larger-text mode, color-blind palette for the table grid.",
], size=12)


# 14. Conclusion / Q&A
s = add_blank()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, INK)
add_rect(s, 0, Inches(3.7), SLIDE_W, Inches(0.06), BRAND)
add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(1.5),
         "Thank you", size=64, bold=True, color=LIGHT)
add_text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(0.6),
         "Smart Restaurant POS  ·  C++17  ·  Dear ImGui  ·  SFML 3.1",
         size=20, color=BRAND)
add_paragraphs(s, Inches(0.7), Inches(4.2), Inches(12), Inches(2), [
    "Live demo:  double-click  dist/smart_pos.exe",
    "Default login:   admin / admin123",
    "Source + planning docs:   ./src/   ./planning/",
], size=18, bullet=False)
add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
         "Questions?", size=24, bold=True, color=LIGHT)


# ============================================================================
out_path = ROOT / "presentation" / "RestaurantPOS_Demo.pptx"
prs.save(str(out_path))
print(f"wrote {out_path}  ({out_path.stat().st_size:,} bytes, {len(prs.slides)} slides)")

# Capture checklist (only print if any placeholders were used).
expected_shots = ["login", "dashboard", "menu", "order", "billing",
                  "inventory", "tables", "kitchen", "analytics"]
missing = [n for n in expected_shots if not (SHOTS_DIR / f"{n}.png").exists()]
if missing:
    print()
    print("Screenshot capture checklist (drop PNGs into presentation/screenshots/):")
    for n in missing:
        print(f"  - {n}.png")
